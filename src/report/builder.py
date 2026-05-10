"""Construção dos entregáveis finais (Etapa 4).

Estratégia adotada (escolhas por compatibilidade no ambiente Windows + numpy 2):
- Gráficos: Plotly + kaleido 0.2.1 → PNGs (alta resolução)
- Tabelas: pandas Styler → HTML embutido
- Relatório técnico (PDF): reportlab (PDF programático, sem GTK/Pango do
  WeasyPrint que costuma quebrar em Windows)
- Sumário executivo (PDF): reportlab, layout condensado
- Apresentação (PPTX): python-pptx
"""

from __future__ import annotations

import logging
from datetime import datetime
from pathlib import Path

import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer,
    Table, TableStyle,
)

from src import paths

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO,
                    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
                    datefmt="%H:%M:%S")


# ---------------------------------------------------------------------------
# Identidade visual (cores reportlab)
# ---------------------------------------------------------------------------

COR_PRIMARIA = colors.HexColor("#1B3A57")
COR_SECUNDARIA = colors.HexColor("#D9A441")
COR_NEUTRA = colors.HexColor("#6B7280")
COR_DESTAQUE = colors.HexColor("#B0413E")
COR_CLARO = colors.HexColor("#E5E7EB")


def _estilos() -> dict[str, ParagraphStyle]:
    base = getSampleStyleSheet()
    out = {
        "h1": ParagraphStyle("h1", parent=base["Heading1"],
                             fontSize=22, textColor=COR_PRIMARIA,
                             spaceAfter=18, leading=26),
        "h2": ParagraphStyle("h2", parent=base["Heading2"],
                             fontSize=16, textColor=COR_PRIMARIA,
                             spaceBefore=18, spaceAfter=10, leading=20),
        "h3": ParagraphStyle("h3", parent=base["Heading3"],
                             fontSize=12, textColor=COR_NEUTRA,
                             spaceBefore=10, spaceAfter=6),
        "p": ParagraphStyle("p", parent=base["Normal"],
                            fontSize=10.5, leading=15, spaceAfter=8,
                            alignment=4),  # 4 = TA_JUSTIFY
        "capa_titulo": ParagraphStyle("capa_titulo", parent=base["Title"],
                                      fontSize=28, textColor=COR_PRIMARIA,
                                      alignment=1, leading=34, spaceAfter=10),
        "capa_sub": ParagraphStyle("capa_sub", parent=base["Normal"],
                                   fontSize=14, textColor=COR_NEUTRA,
                                   alignment=1, spaceAfter=6),
        "destaque_box": ParagraphStyle("destaque_box", parent=base["Normal"],
                                       fontSize=11, leading=16,
                                       textColor=COR_PRIMARIA,
                                       borderColor=COR_SECUNDARIA,
                                       borderWidth=1, borderPadding=10,
                                       backColor=colors.HexColor("#FEF3C7")),
        "rodape": ParagraphStyle("rodape", parent=base["Normal"],
                                 fontSize=8, textColor=COR_NEUTRA,
                                 alignment=1),
    }
    return out


def _capa(es, titulo: str, subtitulo: str) -> list:
    return [
        Spacer(1, 5 * cm),
        Paragraph(titulo, es["capa_titulo"]),
        Spacer(1, 0.5 * cm),
        Paragraph(subtitulo, es["capa_sub"]),
        Spacer(1, 8 * cm),
        Paragraph(f"FGV NPII · ABEVD", es["capa_sub"]),
        Paragraph(f"Maio de 2026", es["capa_sub"]),
        PageBreak(),
    ]


def _df_para_table(df: pd.DataFrame, larguras: list[float] | None = None) -> Table:
    """Converte DataFrame em Table do reportlab."""
    data = [list(df.columns)] + df.values.tolist()
    t = Table(data, colWidths=larguras, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), COR_PRIMARIA),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F9FAFB")]),
        ("GRID", (0, 0), (-1, -1), 0.25, COR_CLARO),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("ALIGN", (1, 1), (-1, -1), "RIGHT"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    return t


def _img(path: Path, largura_cm: float = 16) -> Image:
    img = Image(str(path), width=largura_cm * cm, height=None)
    # mantém aspect ratio
    aspect = img.imageHeight / float(img.imageWidth)
    img.drawHeight = largura_cm * cm * aspect
    return img


# ---------------------------------------------------------------------------
# Carregamento dos artefatos das etapas anteriores
# ---------------------------------------------------------------------------

def _carregar_artefatos() -> dict:
    serie = pd.read_parquet(paths.OUT_TABELAS / "etapa2" / "nano_serie_temporal.parquet")
    ult = serie[serie["periodo"] == "2025T4"]
    recorte = pd.read_parquet(paths.OUT_TABELAS / "etapa3" / "perfil_recorte_abevd.parquet")
    return {
        "serie": serie,
        "ult": ult,
        "recorte": recorte,
        "graficos": paths.OUT_GRAFICOS / "etapa4",
    }


# ---------------------------------------------------------------------------
# Relatório técnico
# ---------------------------------------------------------------------------

def gerar_relatorio_tecnico(versao: str = "v1") -> Path:
    art = _carregar_artefatos()
    es = _estilos()
    out_pdf = paths.OUT_RELATORIOS / f"relatorio-tecnico-{versao}.pdf"
    out_pdf.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(
        str(out_pdf), pagesize=A4,
        leftMargin=2.5 * cm, rightMargin=2.5 * cm,
        topMargin=2 * cm, bottomMargin=2 * cm,
        title="Nano-empreendedores no Brasil — Relatório Técnico",
        author="FGV NPII",
    )
    el: list = []
    el.extend(_capa(
        es,
        "Estimativa e Caracterização do Universo de Nano-empreendedores no Brasil",
        "Relatório Técnico — Série temporal 2025",
    ))

    # 1. Introdução
    el.append(Paragraph("1. Introdução", es["h1"]))
    el.append(Paragraph(
        "Este relatório apresenta a estimativa e caracterização do universo de "
        "<b>nano-empreendedores no Brasil</b> — definidos como trabalhadores por "
        "conta própria com renda anual igual ou inferior a R$ 40 mil — com base "
        "em microdados oficiais da PNAD Contínua (IBGE) e do Cadastro Nacional "
        "de Microempreendedores Individuais (Receita Federal). O estudo abrange "
        "todos os quatro trimestres de 2025, permitindo análise de série temporal "
        "e identificação de tendências.",
        es["p"],
    ))
    el.append(Paragraph(
        "Os resultados servem de subsídio para o planejamento estratégico da "
        "ABEVD, em especial no fortalecimento e expansão da rede de consultoras "
        "de venda direta e no planejamento tributário de longo prazo.",
        es["p"],
    ))

    # 2. Metodologia
    el.append(Paragraph("2. Metodologia", es["h1"]))
    el.append(Paragraph(
        "<b>Fontes (microdados, conforme ADR-004):</b> PNAD Contínua trimestral 2025 "
        "(T1-T4) e Cadastro Nacional MEI snapshots fim do trimestre (2025-03, 06, 09, 12). "
        "<b>Definição operacional:</b> conta-própria (PNADC VD4009=09) com renda "
        "mensal habitual × 12 ≤ R$ 40.000. <b>Expansão amostral:</b> estimador de "
        "Taylor (samplics) com pesos calibrados V1028, estratificado por Estrato "
        "e clusterizado por UPA. IC 95%. <b>Censo 2022:</b> microdados da Amostra "
        "ainda não publicados pelo IBGE (adiados sem nova data) — análise municipal "
        "fica para iteração futura. Detalhes completos em <i>docs/metodologia.md</i>.",
        es["p"],
    ))

    # 3. Validação cruzada
    el.append(Paragraph("3. Validação cruzada", es["h1"]))
    el.append(Paragraph(
        "Antes da apresentação dos resultados, todas as estimativas foram "
        "confrontadas com publicações oficiais (release IBGE PNADC 2025T4, "
        "Painel Sebrae do MEI). 11 métricas testadas, 10 verdes (Δ < 2%) + "
        "1 amarela (validação parcial). Em particular, percentuais de conta-"
        "própria por UF batem <b>exatamente até a primeira decimal</b> com a "
        "publicação oficial do IBGE — confirmando aplicação correta dos pesos "
        "amostrais. Memorial detalhado em <i>docs/validacao-cruzada.md</i>.",
        es["p"],
    ))

    el.append(PageBreak())

    # 4. Estimativa do universo
    el.append(Paragraph("4. Estimativa do universo", es["h1"]))
    el.append(Paragraph("4.1 Brasil — série 2025", es["h2"]))
    g2 = art["graficos"] / "g2_evolucao_brasil.png"
    if g2.exists():
        el.append(_img(g2, 16))
    el.append(Paragraph(
        "<b>Achado principal:</b> o universo nano-empreendedor manteve-se "
        "estável em <b>~19,4 milhões</b> ao longo de 2025, enquanto o cadastro "
        "MEI cresceu cerca de 10-15% no mesmo período em todas as principais "
        "UFs. Isso evidencia uma <b>migração da informalidade para a "
        "formalização</b> em curso — sem expansão do contingente total, mas "
        "com mudança qualitativa na composição.",
        es["p"],
    ))

    el.append(Paragraph("4.2 Distribuição estadual — 2025T4", es["h2"]))
    g1 = art["graficos"] / "g1_nano_por_uf.png"
    if g1.exists():
        el.append(_img(g1, 16))

    el.append(Paragraph("4.3 Cruzamento PNADC × MEI por UF", es["h2"]))
    g3 = art["graficos"] / "g3_pareadas_nano_mei.png"
    if g3.exists():
        el.append(_img(g3, 16))
    el.append(Paragraph(
        "São Paulo apresenta a maior aproximação entre estoque MEI e universo "
        "nano (taxa de formalização aparente de 98%) — porém, este número está "
        "<b>superestimado</b>, pois o teto MEI vigente (R$ 81 mil/ano) é maior "
        "que o teto operacional do estudo (R$ 40 mil), incluindo MEIs que "
        "estão fora do recorte. Em estados como Bahia (37%) e Pará (23%), a "
        "distância entre cadastro e universo aponta espaço relevante para "
        "políticas de formalização.",
        es["p"],
    ))

    el.append(PageBreak())

    # 5. Caracterização socioeconômica
    el.append(Paragraph("5. Caracterização socioeconômica", es["h1"]))
    el.append(Paragraph("5.1 Distribuição etária", es["h2"]))
    g5 = art["graficos"] / "g5_piramide_etaria.png"
    if g5.exists():
        el.append(_img(g5, 14))

    el.append(Paragraph("5.2 Distribuição setorial", es["h2"]))
    g6 = art["graficos"] / "g6_setores_cnae.png"
    if g6.exists():
        el.append(_img(g6, 16))
    el.append(Paragraph(
        "A composição setorial dos nano-empreendedores é dominada por "
        "<b>Agricultura/pecuária (A)</b>, refletindo a alta presença de "
        "trabalhadores rurais autônomos. Em segundo plano figuram "
        "<b>Construção (F)</b>, <b>Outras atividades de serviços (S)</b> "
        "— que inclui salões de beleza, estética, manutenção doméstica — e "
        "<b>Comércio (G)</b>. As seções G+S, em conjunto, representam o "
        "público naturalmente alinhado ao modelo de venda direta da ABEVD.",
        es["p"],
    ))

    el.append(PageBreak())

    # 6. Recorte estratégico ABEVD
    el.append(Paragraph("6. Recorte estratégico para a ABEVD", es["h1"]))
    g4 = art["graficos"] / "g4_recorte_abevd.png"
    if g4.exists():
        el.append(_img(g4, 16))
    el.append(Paragraph(
        "O recorte de <b>mulheres entre 25 e 49 anos atuando em Comércio (G) "
        "ou Outras atividades de serviços (S) como nano-empreendedoras</b> "
        "estima um universo de <b>1,14 milhão de pessoas no Brasil em 2025T4</b>. "
        "São Paulo (247 mil), Minas Gerais (148 mil) e Rio de Janeiro (124 mil) "
        "concentram quase metade desse público, mas o ranking absoluto não deve "
        "ser o único critério: estados como Goiás, Pernambuco e Bahia "
        "apresentam densidade relevante e potencial de expansão proporcional.",
        es["p"],
    ))

    # 7. Limitações
    el.append(Paragraph("7. Limitações", es["h1"]))
    el.append(Paragraph(
        "<b>(a)</b> A taxa de formalização aproximada está superestimada nas UFs "
        "onde o estoque MEI > universo nano estimado, pois o teto MEI (R$ 81 "
        "mil) é maior que o teto operacional (R$ 40 mil); refinamento previsto "
        "em iteração futura. <b>(b)</b> Microdados da Amostra do Censo 2022 "
        "não foram publicados pelo IBGE até a data deste relatório, "
        "impossibilitando análise municipal (ADR-007). <b>(c)</b> Cerca de "
        "14% dos nano-empreendedores aparecem com seção CNAE não declarada "
        "(V4013 ausente). <b>(d)</b> Validação MEI por UF foi parcial — "
        "Painel Sebrae não publica MEI absoluto por UF em formato consultável "
        "via web.",
        es["p"],
    ))

    # 8. Recomendações
    el.append(Paragraph("8. Recomendações para a ABEVD", es["h1"]))
    el.append(Paragraph(
        "<b>1. Priorizar UFs com densidade absoluta + tendência de crescimento "
        "MEI</b>: SP, MG, RJ concentram volume mas estão próximas da saturação "
        "MEI. Estados com formalização ainda baixa (BA, PA, MA, NE em geral) "
        "oferecem terreno fértil para crescimento da rede.",
        es["p"],
    ))
    el.append(Paragraph(
        "<b>2. Apoiar o processo de formalização</b>: a tendência observada "
        "(MEI crescendo ~10-15% ao ano sem expansão do universo nano total) "
        "indica que o público-alvo está naturalmente migrando para a "
        "formalidade. A ABEVD pode acelerar esse movimento via parcerias com "
        "Sebrae e capacitação tributária.",
        es["p"],
    ))
    el.append(Paragraph(
        "<b>3. Recortar comunicação por perfil</b>: o universo é heterogêneo "
        "(rurais no Norte/Nordeste vs urbanos no Sudeste/Sul). Mensagens "
        "uniformes perdem efetividade — segmentação UF × setor CNAE é "
        "viável a partir das tabelas em <i>outputs/tabelas/</i>.",
        es["p"],
    ))

    doc.build(el)
    logger.info("Relatório técnico gerado: %s", out_pdf)
    return out_pdf


# ---------------------------------------------------------------------------
# Sumário executivo
# ---------------------------------------------------------------------------

def gerar_sumario_executivo(versao: str = "v1") -> Path:
    art = _carregar_artefatos()
    es = _estilos()
    out_pdf = paths.OUT_RELATORIOS / f"sumario-executivo-{versao}.pdf"
    out_pdf.parent.mkdir(parents=True, exist_ok=True)

    doc = SimpleDocTemplate(
        str(out_pdf), pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=2 * cm, bottomMargin=2 * cm,
        title="Sumário Executivo — Nano-empreendedores",
        author="FGV NPII",
    )
    el: list = []

    el.append(Paragraph("Sumário Executivo", es["h1"]))
    el.append(Paragraph(
        "Estimativa e Caracterização do Universo de Nano-empreendedores no Brasil — Série 2025",
        es["h3"]
    ))
    el.append(Spacer(1, 0.4 * cm))

    el.append(Paragraph("Principais achados", es["h2"]))
    el.append(Paragraph(
        "• <b>19,16 milhões</b> de nano-empreendedores no Brasil em 2025T4 — "
        "estável ao longo do ano (entre 19,15 M e 19,45 M).<br/>"
        "• <b>13,27 milhões</b> de MEI ativos em dezembro/2025 — crescimento de "
        "10-15% ao ano nas principais UFs.<br/>"
        "• <b>Migração informal → formal</b> em curso: universo total estável + "
        "MEI crescendo significa mudança de composição, não de tamanho.<br/>"
        "• <b>São Paulo lidera</b> em volume absoluto (3,75 M nano; 3,69 M MEI) "
        "e está próximo da saturação. Estados do Norte e Nordeste ainda têm "
        "alta informalidade — espaço para crescimento.<br/>"
        "• <b>Recorte estratégico ABEVD</b>: 1,14 milhão de mulheres 25-49 em "
        "Comércio + Serviços pessoais como nano-empreendedoras — público "
        "naturalmente alinhado ao modelo de venda direta.",
        es["p"],
    ))
    el.append(Spacer(1, 0.3 * cm))

    g2 = art["graficos"] / "g2_evolucao_brasil.png"
    if g2.exists():
        el.append(_img(g2, 16))

    el.append(PageBreak())

    el.append(Paragraph("Recortes estratégicos para a ABEVD", es["h2"]))
    g4 = art["graficos"] / "g4_recorte_abevd.png"
    if g4.exists():
        el.append(_img(g4, 16))

    el.append(Paragraph("Recomendações prioritárias", es["h2"]))
    el.append(Paragraph(
        "<b>1.</b> Priorizar SP, MG, RJ por volume; investir em BA, PA, MA, "
        "NE em geral pelo potencial de crescimento.<br/><br/>"
        "<b>2.</b> Apoiar formalização (parcerias Sebrae) — a tendência "
        "natural já está a favor.<br/><br/>"
        "<b>3.</b> Segmentar comunicação por perfil regional + setorial; "
        "evitar mensagens uniformes.",
        es["destaque_box"],
    ))
    el.append(Spacer(1, 0.5 * cm))

    el.append(Paragraph(
        f"Relatório técnico completo, memoriais de cálculo e validação cruzada "
        f"em <i>docs/</i> e <i>outputs/tabelas/</i>. Documento gerado em "
        f"{datetime.now().strftime('%d/%m/%Y')} pelo FGV NPII para a ABEVD.",
        es["rodape"],
    ))

    doc.build(el)
    logger.info("Sumário executivo gerado: %s", out_pdf)
    return out_pdf


# ---------------------------------------------------------------------------
# Apresentação executiva (PPTX)
# ---------------------------------------------------------------------------

def gerar_apresentacao(versao: str = "v1") -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    art = _carregar_artefatos()
    out_pptx = paths.OUT_RELATORIOS / f"apresentacao-{versao}.pptx"
    out_pptx.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    AZUL = RGBColor(0x1B, 0x3A, 0x57)
    DOURADO = RGBColor(0xD9, 0xA4, 0x41)
    CINZA = RGBColor(0x6B, 0x72, 0x80)

    def slide_titulo(titulo: str, subtitulo: str = ""):
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        tx = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(1.5))
        p = tx.text_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = AZUL
        if subtitulo:
            tx2 = s.shapes.add_textbox(Inches(0.6), Inches(4), Inches(12), Inches(1))
            p2 = tx2.text_frame.paragraphs[0]
            p2.text = subtitulo
            p2.font.size = Pt(20)
            p2.font.color.rgb = CINZA
        return s

    def slide_bullet(titulo: str, bullets: list[str], img: Path | None = None):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        tx_t = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        p = tx_t.text_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = AZUL
        # Bullets à esquerda, imagem à direita (se houver)
        bullet_w = Inches(6) if img else Inches(12)
        tx_b = s.shapes.add_textbox(Inches(0.5), Inches(1.3), bullet_w, Inches(5.5))
        for i, b in enumerate(bullets):
            p = tx_b.text_frame.paragraphs[0] if i == 0 else tx_b.text_frame.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(18)
            p.font.color.rgb = AZUL
            p.space_after = Pt(10)
        if img and img.exists():
            s.shapes.add_picture(str(img), Inches(7), Inches(1.3),
                                 width=Inches(6), height=Inches(5.5))
        return s

    # Slide 1: Capa
    slide_titulo(
        "Nano-empreendedores no Brasil",
        "Estimativa e Caracterização — Série 2025  ·  FGV NPII / ABEVD",
    )

    # Slide 2: Achados principais
    slide_bullet(
        "Achados principais — Brasil 2025",
        [
            "19,16 milhões de nano-empreendedores (renda anual ≤ R$ 40 mil)",
            "Universo estável ao longo de 2025 (entre 19,15 M e 19,45 M)",
            "13,27 milhões de MEI ativos (cadastro, dez/2025)",
            "MEI cresce 10-15% ao ano: migração informal → formal em curso",
            "Validação cruzada com IBGE: 10 de 11 métricas dentro de Δ < 2%",
        ],
        img=art["graficos"] / "g2_evolucao_brasil.png",
    )

    # Slide 3: Top UFs
    slide_bullet(
        "Distribuição estadual",
        [
            "SP, MG, RJ concentram ~35% do universo nano",
            "MA, PA, AM têm maior % de conta-própria sobre ocupados (≥30%)",
            "DF tem menor % conta-própria (17%) — mercado mais formal",
            "Densidade absoluta vs proporcional sugerem estratégias distintas",
        ],
        img=art["graficos"] / "g1_nano_por_uf.png",
    )

    # Slide 4: Cruzamento nano × MEI
    slide_bullet(
        "Estoque MEI vs universo nano",
        [
            "SP próxima da saturação MEI (98% aparente — superestimado)",
            "BA: 37% formalizado · PA: 23% — espaço para crescimento",
            "Cruzamento agregado por UF (não por CPF — restrição administrativa)",
        ],
        img=art["graficos"] / "g3_pareadas_nano_mei.png",
    )

    # Slide 5: Setor CNAE
    slide_bullet(
        "Composição setorial",
        [
            "Agricultura (A) lidera — autônomos rurais",
            "Construção (F), Serviços pessoais (S), Comércio (G), Transformação (C)",
            "Seções G+S = público alinhado ao modelo de venda direta ABEVD",
        ],
        img=art["graficos"] / "g6_setores_cnae.png",
    )

    # Slide 6: Recorte ABEVD
    slide_bullet(
        "Recorte estratégico ABEVD",
        [
            "1,14 milhão de mulheres 25-49 em Comércio + Serviços pessoais (nano)",
            "SP (247k) · MG (148k) · RJ (124k) — quase metade do público",
            "Representa 5-8% do universo nano em cada UF",
            "Público naturalmente alinhado ao modelo de consultoras",
        ],
        img=art["graficos"] / "g4_recorte_abevd.png",
    )

    # Slide 7: Recomendações
    slide_bullet(
        "Recomendações para a ABEVD",
        [
            "1. Priorizar SP/MG/RJ por volume; expandir em BA/PA/MA pelo potencial",
            "2. Apoiar formalização (parcerias Sebrae) — vento já a favor",
            "3. Segmentar comunicação por perfil regional + setorial",
            "4. Reaplicar estudo trimestralmente para acompanhar tendência",
        ],
    )

    # Slide 8: Limitações + próximos passos
    slide_bullet(
        "Limitações e próximos passos",
        [
            "Taxa de formalização superestimada onde teto MEI > teto nano",
            "Censo 2022 (Amostra) não publicado — análise municipal pendente",
            "14% dos nano sem CNAE declarado (V4013 vazia na PNADC)",
            "Próxima iteração: refinar formalização e adicionar trim-2026 quando publicado",
        ],
    )

    prs.save(str(out_pptx))
    logger.info("Apresentação gerada: %s", out_pptx)
    return out_pptx


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def build_all(versao: str = "v1") -> dict[str, Path]:
    paths.ensure_dirs()
    return {
        "relatorio": gerar_relatorio_tecnico(versao),
        "sumario": gerar_sumario_executivo(versao),
        "apresentacao": gerar_apresentacao(versao),
    }


def main() -> None:
    import argparse
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--versao", default="v1")
    args = parser.parse_args()
    out = build_all(args.versao)
    print("Entregáveis gerados:")
    for k, p in out.items():
        print(f"  {k:13s} -> {p}  ({p.stat().st_size/1024:.0f} KB)")


if __name__ == "__main__":
    main()
